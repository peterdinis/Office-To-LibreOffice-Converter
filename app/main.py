"""
FastAPI Office to LibreOffice Converter - Python Libraries Only
Converts Microsoft Office files to LibreOffice formats using Python libraries.
"""

import logging
from io import BytesIO
from collections import OrderedDict
from datetime import datetime, timedelta
from typing import Dict

from fastapi import FastAPI, UploadFile, File, HTTPException, Request
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

# Document processing libraries
from docx import Document
from odf.opendocument import OpenDocumentText, OpenDocumentPresentation
from odf.text import P
from odf.draw import Page, Frame, TextBox
import openpyxl
from pyexcel_ods import save_data
from pptx import Presentation

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Rate limiting storage
rate_limit_storage: Dict[str, list] = {}

# FastAPI application
app = FastAPI(
    title="Office to LibreOffice Converter",
    description="Convert Microsoft Office files to LibreOffice formats using Python libraries",
    version="2.1.0",
)

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["*"],
)

# Rate limiting configuration
RATE_LIMIT_REQUESTS = 10
RATE_LIMIT_WINDOW = 60

# Supported formats (Python libraries only)
SUPPORTED_FORMATS = {
    "excel": ["xlsx", "xls", "xlsm"],
    "word": ["docx"],
    "powerpoint": ["pptx"]
}


def check_rate_limit(client_ip: str) -> bool:
    """Check if client has exceeded rate limit."""
    now = datetime.now()
    
    if client_ip not in rate_limit_storage:
        rate_limit_storage[client_ip] = []
    
    # Remove old requests
    rate_limit_storage[client_ip] = [
        timestamp for timestamp in rate_limit_storage[client_ip]
        if (now - timestamp).total_seconds() < RATE_LIMIT_WINDOW
    ]
    
    if len(rate_limit_storage[client_ip]) >= RATE_LIMIT_REQUESTS:
        return False
    
    rate_limit_storage[client_ip].append(now)
    return True


def get_client_ip(request: Request) -> str:
    """Extract client IP address from request."""
    forwarded_for = request.headers.get("X-Forwarded-For")
    if forwarded_for:
        return forwarded_for.split(",")[0].strip()
    
    real_ip = request.headers.get("X-Real-IP")
    if real_ip:
        return real_ip
    
    return request.client.host


@app.get("/")
async def root():
    """Root endpoint with API information."""
    return {
        "message": "Office to LibreOffice Converter API",
        "version": "2.1.0",
        "supported_formats": SUPPORTED_FORMATS,
        "endpoints": {
            "convert": "/convert/",
            "docs": "/docs"
        }
    }


@app.get("/status/")
async def status():
    """Health check endpoint."""
    return {
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "rate_limit": {
            "requests_per_minute": RATE_LIMIT_REQUESTS,
            "window_seconds": RATE_LIMIT_WINDOW
        }
    }


@app.post("/convert/")
async def convert(request: Request, file: UploadFile = File(...)):
    """Convert Office files to LibreOffice formats."""
    
    # Rate limiting
    client_ip = get_client_ip(request)
    if not check_rate_limit(client_ip):
        logger.warning(f"Rate limit exceeded for IP: {client_ip}")
        raise HTTPException(
            status_code=429, 
            detail=f"Rate limit exceeded. Maximum {RATE_LIMIT_REQUESTS} requests per {RATE_LIMIT_WINDOW} seconds.",
            headers={"Retry-After": str(RATE_LIMIT_WINDOW)}
        )
    
    logger.info(f"Processing conversion request from IP: {client_ip}")
    
    # Validate file
    if not file.filename or "." not in file.filename:
        raise HTTPException(status_code=400, detail="File must have an extension")

    name, ext = file.filename.rsplit(".", 1)
    ext = ext.lower()

    # Read file content
    try:
        contents = await file.read()
        if not contents:
            raise HTTPException(status_code=400, detail="Uploaded file is empty")
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to read file: {str(e)}")

    output_stream = BytesIO()
    filename = ""

    try:
        if ext in SUPPORTED_FORMATS["excel"]:
            # Convert Excel to ODS
            logger.info(f"Converting Excel file {file.filename}")
            
            wb = openpyxl.load_workbook(BytesIO(contents))
            sheet = wb.active
            
            data = OrderedDict()
            data["Sheet1"] = [list(row) for row in sheet.iter_rows(values_only=True)]
            
            save_data(output_stream, data)
            filename = f"{name}.ods"

        elif ext in SUPPORTED_FORMATS["word"]:
            # Convert Word to ODT
            logger.info(f"Converting Word file {file.filename}")
            
            doc = Document(BytesIO(contents))
            odt = OpenDocumentText()
            
            for para in doc.paragraphs:
                if para.text.strip():  # Only add non-empty paragraphs
                    odt.text.addElement(P(text=para.text))
            
            odt.save(output_stream)
            filename = f"{name}.odt"

        elif ext in SUPPORTED_FORMATS["powerpoint"]:
            # Convert PowerPoint to ODP
            logger.info(f"Converting PowerPoint file {file.filename}")
            
            try:
                prs = Presentation(BytesIO(contents))
                odp = OpenDocumentPresentation()
                
                logger.info(f"Processing {len(prs.slides)} slides")
                
                for slide_idx, slide in enumerate(prs.slides):
                    page = Page()
                    
                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                            try:
                                text = shape.text.strip() if shape.text else ""
                                if text:
                                    frame = Frame()
                                    textbox = TextBox()
                                    textbox.addElement(P(text=text))
                                    frame.addElement(textbox)
                                    page.addElement(frame)
                            except Exception:
                                continue
                    
                    odp.presentation.addElement(page)
                
                odp.save(output_stream)
                filename = f"{name}.odp"
                
            except Exception as e:
                logger.error(f"PowerPoint conversion error: {e}")
                raise HTTPException(status_code=500, detail=f"PowerPoint conversion failed: {str(e)}")

        else:
            # Unsupported format
            supported_list = []
            for format_type, extensions in SUPPORTED_FORMATS.items():
                supported_list.extend(extensions)
            
            return JSONResponse(
                status_code=400, 
                content={
                    "error": f"Unsupported file format: {ext}",
                    "supported_formats": supported_list
                }
            )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Conversion failed for {file.filename}: {str(e)}")
        raise HTTPException(status_code=500, detail=f"Conversion failed: {str(e)}")

    # Validate output
    output_stream.seek(0, 2)
    content_length = output_stream.tell()
    output_stream.seek(0)
    
    if content_length == 0:
        raise HTTPException(status_code=500, detail="Converted file is empty")

    # Response headers
    headers = {
        "Content-Disposition": f"attachment; filename={filename}",
        "X-Conversion-Status": "success",
        "X-Rate-Limit-Remaining": str(RATE_LIMIT_REQUESTS - len(rate_limit_storage.get(client_ip, []))),
    }

    logger.info(f"Successfully converted {file.filename} to {filename} ({content_length} bytes)")

    return StreamingResponse(
        output_stream,
        media_type="application/octet-stream",
        headers=headers
    )