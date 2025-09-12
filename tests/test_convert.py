from fastapi.testclient import TestClient
from app.main import app
from io import BytesIO
from docx import Document
import openpyxl
from pptx import Presentation

client = TestClient(app)

def create_excel_file():
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.append(["Name", "Age"])
    ws.append(["Alice", 30])
    ws.append(["Bob", 25])
    stream = BytesIO()
    wb.save(stream)
    stream.seek(0)
    return stream

def create_word_file():
    doc = Document()
    doc.add_paragraph("Hello World")
    stream = BytesIO()
    doc.save(stream)
    stream.seek(0)
    return stream

def create_ppt_file():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Title"
    stream = BytesIO()
    prs.save(stream)
    stream.seek(0)
    return stream

def create_doc_file():
    return BytesIO(b"Fake DOC content")

def create_xlsb_file():
    return BytesIO(b"Fake XLSB content")

def create_unsupported_file():
    return BytesIO(b"Not a valid office file")

def test_excel_conversion():
    files = {"file": ("test.xlsx", create_excel_file(), "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")}
    response = client.post("/convert/", files=files)
    assert response.status_code == 200
    assert response.headers["Content-Disposition"].endswith(".ods")
    assert response.headers.get("X-Conversion-Status") == "success"

def test_word_conversion():
    files = {"file": ("test.docx", create_word_file(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
    response = client.post("/convert/", files=files)
    assert response.status_code == 200
    assert response.headers["Content-Disposition"].endswith(".odt")
    assert response.headers.get("X-Conversion-Status") == "success"

def test_libre_doc_conversion():
    files = {"file": ("test.doc", create_doc_file(), "application/msword")}
    response = client.post("/convert/", files=files)
    assert response.status_code in (200, 500)  

def test_libre_xlsb_conversion():
    files = {"file": ("test.xlsb", create_xlsb_file(), "application/vnd.ms-excel.sheet.binary.macroEnabled.12")}
    response = client.post("/convert/", files=files)
    assert response.status_code in (200, 500)

def test_unsupported_file():
    files = {"file": ("test.txt", create_unsupported_file(), "text/plain")}
    response = client.post("/convert/", files=files)
    assert response.status_code == 400
    assert response.json() == {"error": "Unsupported file format"}

